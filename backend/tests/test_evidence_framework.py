from __future__ import annotations

import asyncio
from datetime import datetime, timezone

from pptx import Presentation

from app.agents.base import AgentContext
from app.agents import pipeline
from app.config import Settings
from app.models import (
    AccountReport,
    Claim,
    DeckSpec,
    EvidenceSignal,
    EvidenceSignalType,
    EvidenceSource,
    EvidenceTableRow,
    ExtractedValue,
    FreshnessWindow,
    ReportSection,
    ResearchMode,
    ResearchRun,
    SlideSpec,
    SourceCredibility,
    SourceSnapshot,
    SourceTier,
)
from app.providers.registry import ProviderRegistry
from app.services.exporter import ReportExporter, _clean_text, _insight_inventory, _leadership_brief
from app.services.report_chat import ReportChatService
from app.storage.evidence_store import EvidenceTableStore


def _context(tmp_path) -> AgentContext:
    run = ResearchRun(company_name="BenchmarkCo", mode=ResearchMode.quick, freshness_window=FreshnessWindow.six_months)
    report = AccountReport(
        run_id=run.id,
        company_name=run.company_name,
        mode=run.mode,
        freshness_window=run.freshness_window,
        sections=[
            ReportSection(id=section_id, title=title, summary="", status="partial", confidence_score=0)
            for section_id, title in pipeline.REPORT_SECTION_ORDER
        ],
        claims=[],
        sources=[],
    )
    return AgentContext(
        run=run,
        report=report,
        providers=ProviderRegistry(Settings(data_dir=tmp_path, allow_live_providers=False)),
    )


def test_source_tier_inference_and_allowed_uses():
    annual_report = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/investors/annual-report-2026.pdf",
            title="Annual Report 2026",
            publisher="BenchmarkCo Investor Relations",
            credibility=SourceCredibility.company_page,
            credibility_score=0.5,
        )
    )
    press_release = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/partnership-announcement",
            title="BenchmarkCo announces AI partnership",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.5,
        )
    )
    careers_page = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/careers/jobs/cloud-engineer",
            title="Cloud Engineer Jobs",
            publisher="BenchmarkCo Careers",
            credibility=SourceCredibility.job_or_career_page,
            credibility_score=0.5,
        )
    )

    assert annual_report.source_tier == SourceTier.tier_1_official_financial
    assert "exact_financial_metric" in annual_report.allowed_uses
    assert press_release.source_tier == SourceTier.tier_2_official_company
    assert "announced_partnership" in press_release.allowed_uses
    assert careers_page.source_tier == SourceTier.tier_4_directional_signal
    assert "hiring_signal" in careers_page.allowed_uses


def test_recency_sensitive_sections_exclude_2024_baseline_when_recent_sources_exist(tmp_path):
    context = _context(tmp_path)
    context.report.generated_at = datetime(2026, 5, 13, tzinfo=timezone.utc)
    stale_source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/press-releases/2025/1/company-reports-full-year-results-2024",
            title="BenchmarkCo reports fourth quarter results and full-year results 2024",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.82,
        )
    )
    recent_source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/press-releases/2026/2/company-launches-ai-investment-program",
            title="BenchmarkCo launches AI investment program",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.82,
        )
    )
    context.report.sources.extend([stale_source, recent_source])

    ranked = pipeline._rank_source_ids_for_section(
        context,
        [stale_source.id, recent_source.id],
        "it_spend",
        include_stale_if_needed=False,
    )

    assert ranked == [recent_source.id]
    assert pipeline._source_is_stale_for_section(context, stale_source, "it_spend") is True
    assert pipeline._source_is_stale_for_section(context, stale_source, "financial_trends") is False


def test_stale_only_claims_are_flagged_for_freshness_sensitive_sections(tmp_path):
    context = _context(tmp_path)
    context.report.generated_at = datetime(2026, 5, 13, tzinfo=timezone.utc)
    stale_source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/press-releases/2025/1/company-reports-full-year-results-2024",
            title="BenchmarkCo reports fourth quarter results and full-year results 2024",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.82,
        )
    )
    context.report.sources.append(stale_source)

    claim_id = pipeline._add_claim(
        context,
        "recent_investments",
        "BenchmarkCo has a recent investment signal based only on the 2024 full-year report.",
        "fact",
        [stale_source.id],
        0.82,
    )

    assert claim_id in pipeline._stale_sensitive_claim_ids(context)


def test_add_claim_creates_evidence_signal_and_dedupes(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-investment",
            title="BenchmarkCo announces AI investment program",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)

    first_id = pipeline._add_claim(
        context,
        "recent_investments",
        "BenchmarkCo announced an AI investment program for automation and cloud modernization.",
        "fact",
        [source.id],
        0.86,
    )
    second_id = pipeline._add_claim(
        context,
        "recent_investments",
        "BenchmarkCo announced an AI investment program for automation and cloud modernization.",
        "fact",
        [source.id],
        0.86,
    )

    assert first_id == second_id
    assert len(context.report.claims) == 1
    assert len(context.report.evidence_signals) == 1
    assert context.report.evidence_signals[0].signal_type == EvidenceSignalType.investment
    assert context.report.evidence_signals[0].signal_strength == "directional"


def test_claim_deduplication_updates_evidence_signals(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/cloud-modernization",
            title="BenchmarkCo cloud modernization",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)
    claim_a = Claim(
        id="claim_a",
        section_id="technology_stack",
        text="BenchmarkCo cited cloud modernization as a technology priority.",
        claim_type="fact",
        evidence_source_ids=[source.id],
        confidence_score=0.8,
    )
    claim_b = Claim(
        id="claim_b",
        section_id="technology_stack",
        text="BenchmarkCo cited cloud modernization as a technology priority.",
        claim_type="fact",
        evidence_source_ids=[source.id],
        confidence_score=0.8,
    )
    context.report.claims.extend([claim_a, claim_b])
    section = pipeline._section(context, "technology_stack")
    section.claim_ids.extend([claim_a.id, claim_b.id])
    context.report.evidence_signals.extend(
        [
            EvidenceSignal(
                id="sig_a",
                section_id="technology_stack",
                signal_type=EvidenceSignalType.technology_stack,
                title="Cloud modernization priority",
                detail=claim_a.text,
                signal_strength="directional",
                source_ids=[source.id],
                claim_ids=[claim_a.id],
                confidence_score=0.8,
            ),
            EvidenceSignal(
                id="sig_b",
                section_id="technology_stack",
                signal_type=EvidenceSignalType.technology_stack,
                title="Cloud modernization priority",
                detail=claim_b.text,
                signal_strength="directional",
                source_ids=[source.id],
                claim_ids=[claim_b.id],
                confidence_score=0.8,
            ),
        ]
    )

    pipeline._dedupe_report_claims(context.report)

    assert len(context.report.claims) == 1
    assert section.claim_ids == ["claim_a"]
    assert len(context.report.evidence_signals) == 1
    assert context.report.evidence_signals[0].claim_ids == ["claim_a"]


def test_evidence_payload_includes_source_tier_and_signals(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-platform",
            title="BenchmarkCo launches AI platform",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)
    pipeline._add_claim(
        context,
        "ai_strategy",
        "BenchmarkCo launched an AI platform to support automation and data workflows.",
        "fact",
        [source.id],
        0.84,
    )

    payload = pipeline._evidence_payload(context, [source.id])

    assert payload[0]["source_tier"] == SourceTier.tier_2_official_company.value
    assert "announced_partnership" in payload[0]["allowed_uses"]
    assert payload[0]["related_evidence_signals"][0]["signal_type"] == EvidenceSignalType.ai_strategy.value
    assert payload[0]["related_table_rows"] == []


def test_refresh_evidence_tables_preserves_claims_signals_sources_and_snapshots(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-platform",
            title="BenchmarkCo launches AI platform",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)
    snapshot = SourceSnapshot(source_id=source.id, text_excerpt="BenchmarkCo launched an AI platform with automation, cloud, and data workflow details.")
    context.report.snapshots.append(snapshot)
    context.report.extracted_values.append(
        ExtractedValue(
            label="Revenue",
            value="USD 1.2 billion",
            unit="USD",
            period="FY2026",
            source_id=source.id,
            exact=True,
        )
    )
    claim_id = pipeline._add_claim(
        context,
        "ai_strategy",
        "BenchmarkCo launched an AI platform to support automation and data workflows.",
        "fact",
        [source.id],
        0.84,
    )
    pipeline._section(context, "ai_strategy").claim_ids.append(claim_id)

    pipeline.refresh_evidence_tables(context)

    row_types = {row.row_type for row in context.report.evidence_table_rows}
    table_names = {row.table_name for row in context.report.evidence_table_rows}
    claim_rows = [row for row in context.report.evidence_table_rows if row.row_type == "claim"]
    signal_rows = [row for row in context.report.evidence_table_rows if row.row_type == "signal"]
    snapshot_rows = [row for row in context.report.evidence_table_rows if row.row_type == "snapshot"]

    assert {"source", "snapshot", "claim", "signal", "extracted_value"}.issubset(row_types)
    assert "ai_strategy_signals" in table_names
    assert "financial_metrics" in table_names
    assert claim_rows and claim_rows[0].claim_ids == [claim_id]
    assert signal_rows and signal_rows[0].source_ids == [source.id]
    assert snapshot_rows and "automation" in snapshot_rows[0].detail
    assert len(context.report.research_anchor.table_row_ids) == len(context.report.evidence_table_rows)


def test_evidence_payload_uses_existing_table_rows_without_omitting_signal_context(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/partnership",
            title="BenchmarkCo partnership",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)
    pipeline._add_claim(
        context,
        "partnerships_deals",
        "BenchmarkCo announced a cloud partnership to accelerate enterprise AI adoption.",
        "fact",
        [source.id],
        0.86,
    )
    pipeline.refresh_evidence_tables(context)

    payload = pipeline._evidence_payload(context, [source.id])

    related_rows = payload[0]["related_table_rows"]
    assert related_rows
    assert any(row["table_name"] == "partnership_signals" for row in related_rows)
    assert "evidence_table_counts" in payload[0]


def test_exporter_builds_insight_inventory_from_evidence_tables(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-investment-partnership",
            title="BenchmarkCo announces AI investment and cloud partnership",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.88,
        )
    )
    context.report.sources.append(source)
    for section_id, text in [
        ("recent_investments", "BenchmarkCo announced an exact USD 50 million AI investment for automation platforms."),
        ("partnerships_deals", "BenchmarkCo announced a cloud partnership to accelerate enterprise AI adoption."),
        ("ai_strategy", "BenchmarkCo launched an AI platform to support data workflows and automation."),
        ("account_priorities", "BenchmarkCo is prioritizing finance automation and customer support modernization."),
        ("executives", "BenchmarkCo leadership signals point to the CIO and CFO as likely buying-center stakeholders."),
        ("hcltech_penetration", "Inferred opportunity: pursue an AI-led automation pilot tied to cloud modernization."),
    ]:
        claim_id = pipeline._add_claim(context, section_id, text, "fact" if section_id != "hcltech_penetration" else "recommendation", [source.id], 0.86)
        pipeline._section(context, section_id).claim_ids.append(claim_id)

    pipeline.refresh_evidence_tables(context)

    inventory = _insight_inventory(context.report, limit_per_section=4)
    titles = {section["title"] for section in inventory}

    assert "Insight Inventory - Investment Signals" in titles
    assert "Insight Inventory - Partnerships and Deals" in titles
    assert "Insight Inventory - AI Strategy Moves" in titles
    assert "Insight Inventory - Function Priorities" in titles
    assert "Insight Inventory - Buying Center" in titles
    assert "Insight Inventory - HCLTech Account Moves" in titles
    assert all(section["rows"] for section in inventory)


def test_leadership_brief_prioritizes_account_moves_from_evidence_tables(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-governance",
            title="BenchmarkCo expands AI governance programs",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.86,
        )
    )
    context.report.sources.append(source)
    for section_id, text in [
        ("recent_investments", "BenchmarkCo announced AI investment signals for regulated automation."),
        ("account_priorities", "BenchmarkCo buying-center pressure points include governance, finance controls, and cloud modernization."),
        ("hcltech_penetration", "Inferred opportunity: HCLTech should pursue an AI governance pilot as the first account move."),
    ]:
        claim_id = pipeline._add_claim(context, section_id, text, "recommendation" if section_id == "hcltech_penetration" else "fact", [source.id], 0.84)
        pipeline._section(context, section_id).claim_ids.append(claim_id)
    pipeline.refresh_evidence_tables(context)

    brief = _leadership_brief(context.report)

    assert [column["title"] for column in brief] == ["What changed", "Why it matters", "What HCLTech should do"]
    assert any("HCLTech" in row.title for column in brief for row in column["rows"])


def test_pptx_leadership_details_expose_full_truncated_row_text(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/buying-center",
            title="BenchmarkCo buying center update",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.86,
        )
    )
    context.report.sources.append(source)
    context.report.evidence_table_rows.append(
        EvidenceTableRow(
            table_name="strategic_priorities",
            row_type="claim",
            section_id="account_priorities",
            title="Taken together, the verified leadership and strategy evidence maps the primary buying center to the C...",
            detail=(
                "Taken together, the verified leadership and strategy evidence maps the primary buying center "
                "to the CFO, CIO, procurement, and transformation leaders for the full account motion."
            ),
            source_ids=[source.id],
            confidence_score=0.86,
        )
    )
    context.report.deck_spec = DeckSpec(
        title="BenchmarkCo",
        subtitle="Test",
        brand_tokens={},
        slides=[
            SlideSpec(
                id="executive_readout",
                title="Executive Readout",
                layout="section",
                bullets=["Leadership synthesis"],
                citation_source_ids=[source.id],
            )
        ],
    )

    pptx_path = tmp_path / "leadership_details.pptx"
    ReportExporter().export_pptx(context.report, pptx_path)
    presentation = Presentation(str(pptx_path))
    deck_text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert "Leadership Brief Details - Why it matters" in deck_text
    assert "CFO, CIO, procurement, and transformation leaders" in deck_text


def test_evidence_table_store_persists_queryable_rows(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/partnership",
            title="BenchmarkCo partnership",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.8,
        )
    )
    context.report.sources.append(source)
    claim_id = pipeline._add_claim(
        context,
        "partnerships_deals",
        "BenchmarkCo announced a cloud partnership to accelerate enterprise AI adoption.",
        "fact",
        [source.id],
        0.86,
    )
    pipeline._section(context, "partnerships_deals").claim_ids.append(claim_id)
    pipeline.refresh_evidence_tables(context)

    store = EvidenceTableStore(tmp_path)
    store.replace_report_rows(context.report)

    counts = store.table_counts(context.report.run_id)
    partnership_rows = store.list_rows(context.report.run_id, table_name="partnership_signals")

    assert counts["partnership_signals"] >= 1
    assert partnership_rows
    assert partnership_rows[0]["source_ids"]


def test_duckdb_evidence_store_searches_relevant_rows(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/cloud-ai-partnership",
            title="BenchmarkCo cloud AI partnership",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.86,
        )
    )
    context.report.sources.append(source)
    claim_id = pipeline._add_claim(
        context,
        "partnerships_deals",
        "BenchmarkCo announced a cloud AI partnership to accelerate enterprise automation adoption.",
        "fact",
        [source.id],
        0.88,
    )
    pipeline._section(context, "partnerships_deals").claim_ids.append(claim_id)
    pipeline.refresh_evidence_tables(context)

    store = EvidenceTableStore(tmp_path)
    store.replace_report_rows(context.report)

    rows = store.search_rows(context.report.run_id, "What cloud AI partnerships are relevant?", limit=5)

    assert rows
    assert any("partnership" in row["title"].lower() or "partnership" in row["detail"].lower() for row in rows)
    assert rows[0]["include_in_analysis"] is True
    assert rows[0]["source_ids"]


def test_duckdb_evidence_store_hybrid_search_uses_embeddings_and_sql_analytics(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-platform",
            title="BenchmarkCo launches AI automation platform",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.88,
        )
    )
    context.report.sources.append(source)
    for section_id, text in [
        ("ai_strategy", "BenchmarkCo launched an AI automation platform for enterprise workflow modernization."),
        ("partnerships_deals", "BenchmarkCo signed a cloud alliance to accelerate analytics programs."),
    ]:
        claim_id = pipeline._add_claim(context, section_id, text, "fact", [source.id], 0.86)
        pipeline._section(context, section_id).claim_ids.append(claim_id)
    pipeline.refresh_evidence_tables(context)

    store = EvidenceTableStore(tmp_path)
    store.replace_report_rows(context.report)
    rows_to_embed = store.rows_needing_embeddings(context.report.run_id, model="test-embedding", dimensions=3)
    embeddings = {}
    for row in rows_to_embed:
        text = f"{row['title']} {row['detail']}".lower()
        embeddings[row["row_id"]] = [1.0, 0.0, 0.0] if "automation" in text else [0.0, 1.0, 0.0]
    store.upsert_embeddings(
        context.report.run_id,
        embeddings_by_row_id=embeddings,
        model="test-embedding",
        dimensions=3,
    )

    rows = store.hybrid_search_rows(
        context.report.run_id,
        "workflow productivity modernization",
        query_embedding=[1.0, 0.0, 0.0],
        model="test-embedding",
        dimensions=3,
        limit=5,
    )
    analytics = store.analytics_snapshot(context.report.run_id, "AI automation")

    assert rows
    assert rows[0]["semantic_score"] == 1.0
    assert "automation" in f"{rows[0]['title']} {rows[0]['detail']}".lower()
    assert analytics["table_counts"]
    assert analytics["signal_type_mix"]
    assert analytics["top_confidence_rows"]


def test_report_chat_uses_duckdb_rows_without_live_provider(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/ai-investment",
            title="BenchmarkCo AI investment",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.84,
        )
    )
    context.report.sources.append(source)
    claim_id = pipeline._add_claim(
        context,
        "recent_investments",
        "BenchmarkCo announced an AI investment program for automation and cloud modernization.",
        "fact",
        [source.id],
        0.84,
    )
    pipeline._section(context, "recent_investments").claim_ids.append(claim_id)
    pipeline.refresh_evidence_tables(context)
    context.run.report = context.report

    store = EvidenceTableStore(tmp_path)
    store.replace_report_rows(context.report)
    service = ReportChatService(Settings(data_dir=tmp_path, allow_live_providers=False))
    rows, analytics, retrieval_mode = asyncio.run(
        service.retrieve(
            run=context.run,
            question="AI investment",
            evidence_store=store,
            limit=5,
        )
    )

    response = asyncio.run(
        service.answer(
            run=context.run,
            question="What AI investment signals exist?",
            evidence_rows=rows,
            analytics=analytics,
            retrieval_mode=retrieval_mode,
        )
    )

    assert response.provider == "duckdb_retrieval"
    assert response.retrieval_mode == "duckdb_keyword_sql"
    assert response.analytics["table_counts"]
    assert response.evidence_rows
    assert "DuckDB evidence mart" in response.answer


def test_exporter_clean_text_replaces_internal_source_ids_and_scaffold_text():
    cleaned = _clean_text(
        "Company identity, official-source discovery, and report metadata are configured. [src_123456789abc] cited source",
        {"src_123456789abc": 7},
    )

    assert "S7" in cleaned
    assert "cited source" not in cleaned.lower()
    assert "Company overview remains partial" in cleaned


def test_exporter_clean_text_keeps_currency_context_for_financial_metrics():
    cleaned = _clean_text(
        "Revenue was \u20b91,62,990 crore, large-deal TCV was $9.8 billion, and EU revenue was \u20ac2.5 billion.",
        {},
    )

    assert "INR 1,62,990 crore" in cleaned
    assert "USD 9.8 billion" in cleaned
    assert "EUR 2.5 billion" in cleaned
    assert "\u20b9" not in cleaned
    assert "\u20ac" not in cleaned


def test_hcl_strategy_playbook_is_industry_agnostic_for_non_telecom(tmp_path):
    context = _context(tmp_path)
    source = pipeline._enrich_source(
        EvidenceSource(
            url="https://example.com/newsroom/retail-ai-cloud",
            title="BenchmarkCo expands AI and cloud programs",
            publisher="BenchmarkCo",
            credibility=SourceCredibility.company_page,
            credibility_score=0.82,
        )
    )
    context.report.sources.append(source)
    for section_id, text in [
        ("recent_investments", "BenchmarkCo is investing in AI automation and cloud data modernization for retail operations."),
        ("technology_stack", "BenchmarkCo cited data, analytics, and cloud engineering priorities."),
        ("footprint_hiring", "BenchmarkCo hiring signals emphasize data engineering and AI skills."),
    ]:
        claim_id = pipeline._add_claim(context, section_id, text, "fact", [source.id], 0.82)
        pipeline._section(context, section_id).claim_ids.append(claim_id)

    pipeline._apply_hcl_strategy_playbook(context, [source.id])

    hcl_section = pipeline._section(context, "hcltech_penetration")
    summary = hcl_section.summary.lower()
    assert "ai-led productivity" in summary
    assert "customer data, commerce, and supply-chain" in summary
    assert "ran" not in summary
    assert "oss/bss" not in summary
