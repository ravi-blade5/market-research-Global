from __future__ import annotations

import json
import re
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas

from app.brand import BRAND_COLORS
from app.models import AccountReport, Artifact, Claim, EvidenceSource, EvidenceTableRow, QualityCheck, ReportSection, SlideSpec


SOURCE_ID_RE = re.compile(r"src_[0-9a-f]{8,16}")
SCAFFOLD_TERMS = (
    "schema implemented",
    "hooks are ready",
    "metadata are configured",
    "source discovery is configured",
    "scaffold",
)
CLIENT_TEXT_REPLACEMENTS = {
    "Company identity, official-source discovery, and report metadata are configured.": (
        "Company overview remains partial until official company, investor, or filing evidence is extracted into exact facts."
    ),
    "Priority taxonomy is ready; live source extraction will populate priorities per business function.": (
        "Account priorities are reported only when section-specific filings, earnings calls, press releases, or credible news support them."
    ),
    "Technology stack schema implemented with confidence scoring and source trail requirements.": (
        "Technology stack analysis is partial; no vendor, platform, or tool will be listed without public job, official, partner, or credible evidence."
    ),
    "Apify/Firecrawl-backed job and location extraction hooks are ready.": (
        "Hiring and footprint analysis is partial unless public job, career, or official location evidence provides exact role, skill, and location signals."
    ),
    "Signal validation policy is configured.": (
        "Key signals are reported only when supported by official announcements or reputable news within the selected freshness window."
    ),
    "Buying-center map schema implemented.": (
        "The buying-center map is partial until verified executive and functional-leader evidence is available for named stakeholders."
    ),
    "AI maturity and opportunity framework is implemented.": (
        "AI strategy assessment is pending verified evidence on AI investments, partnerships, offerings, adoption, and roadmap signals."
    ),
    "Opportunity-led account penetration playbook scaffold.": (
        "HCLTech account-penetration guidance is pending verified account priorities, technology, sourcing, executive, and AI-strategy evidence."
    ),
    "Consensus engine ready; live runs will rank account moves by evidence strength and strategic confidence.": (
        "Consensus recommendation is pending verified section-level evidence and should not overstate unsupported account moves."
    ),
    "Source-grounded report scaffold with fail-closed unavailable fields.": (
        "Source-grounded report with unsupported fields marked unavailable."
    ),
    "unavailable until official filings or investor materials are extracted by a live provider": (
        "exact value not found in the extracted official evidence for this run"
    ),
    "Pending extraction of an exact revenue value from annual, quarterly, investor, or filing evidence.": (
        "Exact revenue value not found in the extracted official financial evidence for this run."
    ),
    "Pending extraction of an exact R&D value from annual, quarterly, investor, or filing evidence.": (
        "Exact R&D value not found in the extracted official financial evidence for this run."
    ),
}
TEXT_TRANSLATION = str.maketrans(
    {
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2013": "-",
        "\u2014": "-",
        "\u2022": "-",
        "\u00a0": " ",
    }
)
CURRENCY_PREFIX_PATTERNS = (
    (re.compile(r"(?i)\bUS\$\s*(?=\d)"), "USD "),
    (re.compile(r"(?i)\bA\$\s*(?=\d)"), "AUD "),
    (re.compile(r"(?i)\bC\$\s*(?=\d)"), "CAD "),
    (re.compile(r"(?i)\bS\$\s*(?=\d)"), "SGD "),
    (re.compile(r"\u20b9\s*(?=\d)"), "INR "),
    (re.compile(r"\u20ac\s*(?=\d)"), "EUR "),
    (re.compile(r"\u00a3\s*(?=\d)"), "GBP "),
    (re.compile(r"\u00a5\s*(?=\d)"), "JPY "),
    (re.compile(r"\u20a9\s*(?=\d)"), "KRW "),
    (re.compile(r"(?<![A-Za-z])\$\s*(?=\d)"), "USD "),
)
CURRENCY_CODE_SPACING_RE = re.compile(r"\b(INR|USD|EUR|GBP|JPY|AUD|CAD|SGD|KRW)\s*(?=\d)")


def _rgb(name: str) -> RGBColor:
    r, g, b = BRAND_COLORS[name].rgb
    return RGBColor(r, g, b)


def _hex(name: str) -> str:
    return BRAND_COLORS[name].hex


def _rl_color(name: str) -> colors.HexColor:
    return colors.HexColor(_hex(name))


def _source_numbers(report: AccountReport) -> dict[str, int]:
    numbered_sources = [source for source in report.sources if source.url.startswith("http")]
    return {source.id: index + 1 for index, source in enumerate(numbered_sources)}


def _public_sources(report: AccountReport) -> list[EvidenceSource]:
    return [source for source in report.sources if source.url.startswith("http")]


def _clean_text(text: str, source_numbers: dict[str, int]) -> str:
    def replace_source(match: re.Match[str]) -> str:
        source_id = match.group(0)
        number = source_numbers.get(source_id)
        return f"S{number}" if number else "source"

    cleaned = re.sub(r"\[[^\]]*(src_[0-9a-f]{8,16})[^\]]*\]", lambda match: replace_source(re.search(SOURCE_ID_RE, match.group(0))), text)
    cleaned = SOURCE_ID_RE.sub(replace_source, cleaned)
    cleaned = re.sub(r"\[\s*[,;:\s]*\]", "", cleaned)
    cleaned = re.sub(r"\[\s*(?:cited source[\s,;]*)+\]", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s*cited source(?:\s*,\s*cited source)*", "", cleaned, flags=re.IGNORECASE)
    cleaned = cleaned.replace("R&D;", "R&D").strip()
    for old, new in CLIENT_TEXT_REPLACEMENTS.items():
        cleaned = cleaned.replace(old, new)
    return _safe_text(cleaned)


def _safe_text(text: str) -> str:
    translated = _normalize_currency_markers(text).translate(TEXT_TRANSLATION)
    return translated.encode("latin-1", errors="ignore").decode("latin-1")


def _normalize_currency_markers(text: str) -> str:
    normalized = text
    for pattern, replacement in CURRENCY_PREFIX_PATTERNS:
        normalized = pattern.sub(replacement, normalized)
    return CURRENCY_CODE_SPACING_RE.sub(r"\1 ", normalized)


def _readable_source_title(source: EvidenceSource) -> str:
    title = _safe_text(source.title or "").strip()
    if not title or title.startswith("http") or title == source.url:
        parsed = urlparse(source.url)
        domain = parsed.netloc.replace("www.", "") or "source"
        path_hint = parsed.path.strip("/").replace("-", " ").replace("_", " ")
        title = f"{domain}: {path_hint[:54]}" if path_hint else domain
    return _truncate(title, 86)


def _unique(items: list[str]) -> list[str]:
    seen: set[str] = set()
    unique_items: list[str] = []
    for item in items:
        if item not in seen:
            seen.add(item)
            unique_items.append(item)
    return unique_items


def _claims_for_section(report: AccountReport, section: ReportSection) -> list[Claim]:
    claim_ids = set(section.claim_ids)
    return [claim for claim in report.claims if claim.id in claim_ids]


def _citation_ids_for_section(report: AccountReport, section: ReportSection) -> list[str]:
    source_ids = [
        source_id
        for claim in _claims_for_section(report, section)
        for source_id in claim.evidence_source_ids
    ]
    return _unique(source_ids)


def _section_by_id(report: AccountReport, section_id: str) -> ReportSection | None:
    return next((section for section in report.sections if section.id == section_id), None)


def _status_fill_name(status: str | None) -> str:
    if status == "complete":
        return "tech_blue"
    if status == "unavailable":
        return "tech_gray"
    return "light_blue"


def _truncate(text: str, limit: int) -> str:
    cleaned = " ".join(text.split())
    return cleaned if len(cleaned) <= limit else cleaned[: max(0, limit - 3)].rstrip() + "..."


def _citation_labels(report: AccountReport, source_ids: list[str], source_numbers: dict[str, int], limit: int = 5) -> list[str]:
    sources_by_id = {source.id: source for source in report.sources}
    labels: list[str] = []
    for source_id in source_ids:
        source = sources_by_id.get(source_id)
        if not source or not source.url.startswith("http"):
            continue
        number = source_numbers.get(source.id)
        if not number:
            continue
        publisher = f" - {_safe_text(source.publisher)}" if source.publisher else ""
        labels.append(f"S{number}: {_readable_source_title(source)}{publisher}")
        if len(labels) >= limit:
            break
    return labels


def _summary_bullets(text: str, source_numbers: dict[str, int], limit: int = 4) -> list[str]:
    cleaned = _clean_text(text, source_numbers)
    if not cleaned:
        return ["No source-backed synthesis available."]
    sentences = re.split(r"(?<=[.!?])\s+", cleaned)
    bullets = [sentence.strip() for sentence in sentences if sentence.strip()]
    return bullets[:limit] or [cleaned]


def _wrap_text(
    text: str,
    font_name: str,
    font_size: int,
    max_width: float,
    max_lines: int | None = None,
    add_ellipsis: bool = True,
) -> list[str]:
    words = text.replace("\n", " ").split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = word
        while pdfmetrics.stringWidth(current, font_name, font_size) > max_width and len(current) > 12:
            cut = max(12, int(len(current) * max_width / max(pdfmetrics.stringWidth(current, font_name, font_size), 1)))
            lines.append(current[:cut])
            current = current[cut:]
        if max_lines and len(lines) >= max_lines:
            break
    if current and (not max_lines or len(lines) < max_lines):
        lines.append(current)
    if max_lines and len(lines) > max_lines:
        lines = lines[:max_lines]
    if add_ellipsis and max_lines and len(lines) == max_lines and words:
        last = lines[-1]
        if len(" ".join(words)) > len(" ".join(lines)):
            lines[-1] = last[:-3].rstrip() + "..." if len(last) > 12 else last + "..."
    return lines


def _draw_wrapped(
    page: canvas.Canvas,
    text: str,
    x: float,
    y: float,
    max_width: float,
    font_name: str,
    font_size: int,
    leading: int,
    color: colors.Color,
    max_lines: int | None = None,
    add_ellipsis: bool = True,
) -> float:
    page.setFillColor(color)
    page.setFont(font_name, font_size)
    for line in _wrap_text(text, font_name, font_size, max_width, max_lines, add_ellipsis=add_ellipsis):
        page.drawString(x, y, line)
        y -= leading
    return y


def _claim_status_data(report: AccountReport) -> list[dict[str, int | str]]:
    statuses = ["verified", "unavailable", "rejected", "pending"]
    return [{"label": status.title(), "value": len([claim for claim in report.claims if claim.verification_status == status])} for status in statuses]


def _source_mix_data(report: AccountReport) -> list[dict[str, int | str]]:
    counts: dict[str, int] = {}
    for source in _public_sources(report):
        tier = source.source_tier.value if source.source_tier else source.credibility.value
        label = tier.replace("_", " ").title()
        counts[label] = counts.get(label, 0) + 1
    return [{"label": key, "value": value} for key, value in sorted(counts.items(), key=lambda item: item[1], reverse=True)]


def _signal_mix_data(report: AccountReport) -> list[dict[str, int | str]]:
    counts: dict[str, int] = {}
    for signal in report.evidence_signals:
        label = signal.signal_type.value.replace("_", " ").title()
        counts[label] = counts.get(label, 0) + 1
    return [{"label": key, "value": value} for key, value in sorted(counts.items(), key=lambda item: item[1], reverse=True)]


def _table_mix_data(report: AccountReport) -> list[dict[str, int | str]]:
    counts: dict[str, int] = {}
    for row in report.evidence_table_rows:
        label = row.table_name.replace("_", " ").title()
        counts[label] = counts.get(label, 0) + 1
    return [{"label": key, "value": value} for key, value in sorted(counts.items(), key=lambda item: item[1], reverse=True)]


def _coverage_data(report: AccountReport) -> list[dict[str, int | float | str]]:
    return [
        {
            "label": section.title,
            "value": len(section.claim_ids),
            "confidence": section.confidence_score,
            "status": section.status,
        }
        for section in report.sections
    ]


def _metric_counts(report: AccountReport) -> dict[str, int]:
    return {
        "Sources": len(_public_sources(report)),
        "Snapshots": len([snapshot for snapshot in report.snapshots if snapshot.text_excerpt]),
        "Claims": len(report.claims),
        "Rows": len(report.evidence_table_rows),
    }


INSIGHT_INVENTORY_SPECS = (
    {
        "title": "Insight Inventory - Investment Signals",
        "subtitle": "Recent investments, capital allocation, acquisitions, AI/cloud/data programs, and facility moves",
        "table_names": ("investment_signals", "it_investment_signals"),
        "section_ids": ("recent_investments", "it_spend"),
    },
    {
        "title": "Insight Inventory - Partnerships and Deals",
        "subtitle": "Partnerships, alliances, customer deals, partner announcements, and commercial moves",
        "table_names": ("partnership_signals", "news_signals"),
        "section_ids": ("partnerships_deals", "key_signals"),
    },
    {
        "title": "Insight Inventory - AI Strategy Moves",
        "subtitle": "AI products, AI partnerships, adoption signals, roadmap language, and AI-led investment themes",
        "table_names": ("ai_strategy_signals",),
        "section_ids": ("ai_strategy",),
    },
    {
        "title": "Insight Inventory - Function Priorities",
        "subtitle": "Function-by-function priorities and business workload signals from the evidence graph",
        "table_names": ("strategic_priorities",),
        "section_ids": ("account_priorities",),
    },
    {
        "title": "Insight Inventory - Buying Center",
        "subtitle": "Executives, functional leaders, buying-center implications, and stakeholder entry points",
        "table_names": ("executive_buying_center",),
        "section_ids": ("executives",),
    },
    {
        "title": "Insight Inventory - HCLTech Account Moves",
        "subtitle": "Opportunity hypotheses, pursuit motions, first moves, and consensus account recommendations",
        "table_names": ("opportunity_hypotheses", "consensus_moves"),
        "section_ids": ("hcltech_penetration", "consensus"),
        "prefer_keywords": (
            "opportunity",
            "hcltech",
            "account move",
            "entry",
            "pilot",
            "pursue",
            "first 30",
            "first 90",
            "play",
            "recommend",
            "solution",
            "workstream",
            "factory",
        ),
    },
)


def _row_strength(row: EvidenceTableRow) -> str:
    raw = (
        row.normalized_fields.get("signal_strength")
        or row.normalized_fields.get("claim_type")
        or row.normalized_fields.get("verification_status")
        or row.row_type
    )
    return str(raw).replace("_", " ").title()


def _row_display_text(row: EvidenceTableRow, source_numbers: dict[str, int]) -> str:
    title = _clean_text(row.title, source_numbers)
    detail = _clean_text(row.detail, source_numbers)
    if title.endswith("...") and detail and len(detail) > len(title):
        return detail
    return title


def _source_refs(source_ids: list[str], source_numbers: dict[str, int], limit: int = 4) -> str:
    labels = [f"S{source_numbers[source_id]}" for source_id in source_ids if source_id in source_numbers]
    if not labels:
        return "Evidence pack"
    shown = labels[:limit]
    suffix = f" +{len(labels) - limit}" if len(labels) > limit else ""
    return ", ".join(shown) + suffix


def _insight_row_sort_key(row: EvidenceTableRow) -> tuple[int, int, float, int, str]:
    strength = str(row.normalized_fields.get("signal_strength") or row.normalized_fields.get("claim_type") or "")
    strength_rank = {"exact": 0, "fact": 0, "directional": 1, "inferred": 2, "recommendation": 2}.get(strength, 3)
    row_rank = {"signal": 0, "claim": 1, "extracted_value": 2, "section_summary": 3}.get(row.row_type, 4)
    return (row_rank, strength_rank, -row.confidence_score, -len(row.source_ids), row.title.lower())


def _keyword_hit_count(row: EvidenceTableRow, keywords: tuple[str, ...]) -> int:
    if not keywords:
        return 0
    haystack = f"{row.title} {row.detail}".lower()
    return sum(1 for keyword in keywords if keyword in haystack)


def _insight_rows(
    report: AccountReport,
    table_names: tuple[str, ...],
    section_ids: tuple[str, ...] = (),
    prefer_keywords: tuple[str, ...] = (),
    limit: int = 8,
) -> list[EvidenceTableRow]:
    candidate_rows = [
        row
        for row in report.evidence_table_rows
        if row.table_name in table_names
        and row.include_in_analysis
        and row.row_type in {"signal", "claim", "extracted_value", "section_summary"}
        and (not section_ids or row.section_id in section_ids or row.section_id is None)
    ]
    if section_ids and not candidate_rows:
        candidate_rows = [
            row
            for row in report.evidence_table_rows
            if row.table_name in table_names
            and row.include_in_analysis
            and row.row_type in {"signal", "claim", "extracted_value", "section_summary"}
        ]
    if prefer_keywords:
        preferred_rows = [row for row in candidate_rows if _keyword_hit_count(row, prefer_keywords) > 0]
        if preferred_rows:
            candidate_rows = preferred_rows
    seen: set[tuple[str, str]] = set()
    rows: list[EvidenceTableRow] = []
    for row in sorted(candidate_rows, key=lambda row: (-_keyword_hit_count(row, prefer_keywords), *_insight_row_sort_key(row))):
        key = (_clean_text(row.title, {}), _clean_text(row.detail[:240], {}))
        if key in seen:
            continue
        seen.add(key)
        rows.append(row)
        if len(rows) >= limit:
            break
    return rows


def _insight_inventory(report: AccountReport, limit_per_section: int = 8) -> list[dict[str, object]]:
    inventory: list[dict[str, object]] = []
    for spec in INSIGHT_INVENTORY_SPECS:
        rows = _insight_rows(
            report,
            tuple(spec["table_names"]),
            tuple(spec.get("section_ids", ())),
            tuple(spec.get("prefer_keywords", ())),
            limit=limit_per_section,
        )
        if rows:
            inventory.append({**spec, "rows": rows})
    return inventory


def _leadership_brief(report: AccountReport) -> list[dict[str, object]]:
    investment_rows = _insight_rows(
        report,
        ("investment_signals", "partnership_signals", "ai_strategy_signals"),
        prefer_keywords=("announced", "investment", "partnership", "collaboration", "acquisition", "launch", "ai", "cloud"),
        limit=3,
    )
    buyer_rows = _insight_rows(
        report,
        ("strategic_priorities", "executive_buying_center", "it_investment_signals"),
        prefer_keywords=("chief", "cfo", "cio", "buyer", "buying", "leader", "executive", "priority", "stakeholder", "governance", "modernization"),
        limit=3,
    )
    move_rows = _insight_rows(
        report,
        ("opportunity_hypotheses", "consensus_moves"),
        ("hcltech_penetration", "consensus"),
        ("opportunity", "hcltech", "entry", "pilot", "account move", "first 30", "first 90", "workstream"),
        limit=3,
    )
    return [
        {
            "title": "What changed",
            "subtitle": "Evidence-backed signals that should change the account view",
            "rows": investment_rows,
        },
        {
            "title": "Why it matters",
            "subtitle": "Buyer pressure, functional priority, and buying-center implications",
            "rows": buyer_rows,
        },
        {
            "title": "What HCLTech should do",
            "subtitle": "Highest-confidence entry plays and first account moves",
            "rows": move_rows,
        },
    ]


class ReportExporter:
    def export_all(self, report: AccountReport, pptx_path: Path, pdf_path: Path, evidence_path: Path) -> list[Artifact]:
        pptx_artifact = self.export_pptx(report, pptx_path)
        pdf_artifact = self.export_pdf(report, pdf_path)
        evidence_artifact = self.export_evidence(report, evidence_path)
        artifacts = [pptx_artifact, pdf_artifact, evidence_artifact]
        artifact_checks = [check for artifact in artifacts for check in artifact.quality_checks]
        artifact_check_names = {check.name for check in artifact_checks}
        report.quality_checks = [check for check in report.quality_checks if check.name not in artifact_check_names]
        report.quality_checks.extend(artifact_checks)
        return artifacts

    def export_pptx(self, report: AccountReport, path: Path) -> Artifact:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        source_numbers = _source_numbers(report)

        slides = report.deck_spec.slides if report.deck_spec else []
        inserted_inventory = False
        for slide_spec in slides:
            self._add_slide(prs, slide_spec, report, source_numbers)
            if slide_spec.id == "executive_readout":
                self._add_pptx_leadership_brief(prs, report, source_numbers)
                self._add_pptx_leadership_details(prs, report, source_numbers)
            if slide_spec.id == "section_coverage":
                self._add_pptx_insight_inventory(prs, report, source_numbers)
                inserted_inventory = True
        if not inserted_inventory:
            self._add_pptx_insight_inventory(prs, report, source_numbers)

        prs.save(path)
        return Artifact(
            kind="pptx",
            path=str(path),
            quality_checks=[
                QualityCheck(name="pptx_created", passed=path.exists(), message=f"PPTX written to {path}"),
                QualityCheck(name="brand_palette_used", passed=True, message="Mandatory palette used in title, fills, and accents."),
                QualityCheck(
                    name="pptx_human_citations",
                    passed=True,
                    message="PPTX uses reader-facing source labels instead of internal source IDs.",
                ),
            ],
        )

    def _add_slide(self, prs: Presentation, slide_spec: SlideSpec, report: AccountReport, source_numbers: dict[str, int]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        section = _section_by_id(report, slide_spec.id)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb("dark_blue" if slide_spec.layout == "cover" else "light_blue")

        if slide_spec.layout == "cover":
            portal_tag = slide.shapes.add_textbox(Inches(0.72), Inches(0.52), Inches(5.5), Inches(0.32))
            portal_frame = portal_tag.text_frame
            portal_frame.clear()
            portal_p = portal_frame.paragraphs[0]
            portal_p.text = "HCLTech Market Research Portal"
            portal_p.font.size = Pt(12)
            portal_p.font.bold = True
            portal_p.font.color.rgb = _rgb("light_blue")
            accent = slide.shapes.add_shape(1, Inches(0), Inches(5.75), Inches(13.333), Inches(0.42))
            accent.fill.solid()
            accent.fill.fore_color.rgb = _rgb("tech_purple")
            accent.line.fill.background()
            title_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.3), Inches(11.8), Inches(1.0))
        else:
            page_panel = slide.shapes.add_shape(1, Inches(0.42), Inches(1.05), Inches(12.5), Inches(5.75))
            page_panel.fill.solid()
            page_panel.fill.fore_color.rgb = _rgb("tech_gray")
            page_panel.line.color.rgb = _rgb("mid_blue")
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.82))
            header.fill.solid()
            header.fill.fore_color.rgb = _rgb("dark_blue")
            header.line.fill.background()
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.19), Inches(9.2), Inches(0.45))
            if section:
                status = slide.shapes.add_shape(1, Inches(10.25), Inches(0.18), Inches(1.18), Inches(0.34))
                status.fill.solid()
                status.fill.fore_color.rgb = _rgb(_status_fill_name(section.status))
                status.line.fill.background()
                status_frame = status.text_frame
                status_frame.clear()
                status_p = status_frame.paragraphs[0]
                status_p.text = section.status.upper()
                status_p.font.size = Pt(8)
                status_p.font.bold = True
                status_p.font.color.rgb = _rgb("dark_blue")
                status_p.alignment = PP_ALIGN.CENTER
                evidence_count = len(_citation_ids_for_section(report, section))
                evidence_label = slide.shapes.add_textbox(Inches(11.55), Inches(0.22), Inches(1.3), Inches(0.25))
                evidence_label.text_frame.text = f"{evidence_count} sources"
                evidence_label.text_frame.paragraphs[0].font.size = Pt(8)
                evidence_label.text_frame.paragraphs[0].font.color.rgb = _rgb("light_blue")

        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_spec.title
        title_p.font.size = Pt(24 if slide_spec.layout != "cover" else 34)
        title_p.font.bold = True
        title_p.font.color.rgb = _rgb("light_blue" if slide_spec.layout == "cover" else "tech_gray")

        body_width = Inches(5.55 if slide_spec.chart else 7.75 if slide_spec.layout != "cover" else 11.6)
        body = slide.shapes.add_textbox(Inches(0.72), Inches(1.35 if slide_spec.layout != "cover" else 2.65), body_width, Inches(4.95))
        frame = body.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        frame.clear()

        bullets = [_clean_text(bullet, source_numbers) for bullet in (slide_spec.bullets or ["No content available."])]
        for idx, bullet in enumerate(bullets[:5]):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            p.text = _truncate(bullet, 265 if slide_spec.layout != "cover" else 180)
            p.level = 0
            p.font.size = Pt(17 if slide_spec.layout == "cover" else 12)
            p.font.color.rgb = _rgb("tech_gray" if slide_spec.layout == "cover" else "dark_purple")
            p.space_after = Pt(7)

        if slide_spec.layout == "cover":
            self._add_pptx_metric_cards(slide, report)
        elif slide_spec.chart:
            self._add_pptx_chart(slide, slide_spec.chart)
        else:
            citation_box = slide.shapes.add_shape(1, Inches(9.18), Inches(1.35), Inches(3.4), Inches(4.95))
            citation_box.fill.solid()
            citation_box.fill.fore_color.rgb = _rgb("light_blue")
            citation_box.line.color.rgb = _rgb("mid_blue")
            citation_text = slide.shapes.add_textbox(Inches(9.42), Inches(1.58), Inches(2.92), Inches(4.42))
            citation_frame = citation_text.text_frame
            citation_frame.word_wrap = True
            citation_frame.clear()
            heading = citation_frame.paragraphs[0]
            heading.text = "Evidence Trail"
            heading.font.size = Pt(13)
            heading.font.bold = True
            heading.font.color.rgb = _rgb("dark_blue")
            labels = _citation_labels(report, slide_spec.citation_source_ids, source_numbers, limit=5)
            for label in labels or ["No public source attached."]:
                p = citation_frame.add_paragraph()
                p.text = label
                p.font.size = Pt(8)
                p.font.color.rgb = _rgb("dark_purple")
                p.space_after = Pt(5)

        footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.25))
        footer_frame = footer.text_frame
        footer_frame.clear()
        footer_p = footer_frame.paragraphs[0]
        footer_p.text = "HCLTech Market Research Portal | Source-grounded AI report"
        footer_p.font.size = Pt(8)
        footer_p.font.color.rgb = _rgb("light_blue" if slide_spec.layout == "cover" else "dark_blue")
        footer_p.alignment = PP_ALIGN.RIGHT

    def _add_pptx_metric_cards(self, slide, report: AccountReport) -> None:
        metrics = _metric_counts(report)
        x = 0.72
        for label, value in metrics.items():
            card = slide.shapes.add_shape(1, Inches(x), Inches(4.6), Inches(2.65), Inches(0.86))
            card.fill.solid()
            card.fill.fore_color.rgb = _rgb("light_blue")
            card.line.fill.background()
            tf = card.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = _rgb("dark_blue")
            label_p = tf.add_paragraph()
            label_p.text = label
            label_p.font.size = Pt(9)
            label_p.font.color.rgb = _rgb("dark_purple")
            x += 2.88

    def _add_pptx_chart(self, slide, chart: dict) -> None:
        data = chart.get("data") or []
        title = chart.get("title") or "Evidence Chart"
        chart_box = slide.shapes.add_shape(1, Inches(6.7), Inches(1.18), Inches(5.95), Inches(4.95))
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = _rgb("light_blue")
        chart_box.line.color.rgb = _rgb("mid_blue")
        heading = slide.shapes.add_textbox(Inches(6.95), Inches(1.38), Inches(5.4), Inches(0.3))
        heading.text_frame.text = title
        heading.text_frame.paragraphs[0].font.bold = True
        heading.text_frame.paragraphs[0].font.size = Pt(13)
        heading.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_blue")
        if not data:
            return
        max_value = max(float(item.get("value", 0) or 0) for item in data) or 1
        y = 1.88
        for item in data[:8]:
            label = str(item.get("label", ""))[:34]
            value = float(item.get("value", 0) or 0)
            label_box = slide.shapes.add_textbox(Inches(6.95), Inches(y), Inches(2.0), Inches(0.22))
            label_box.text_frame.text = label
            label_box.text_frame.paragraphs[0].font.size = Pt(7)
            label_box.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_purple")
            bar_width = 2.9 * (value / max_value)
            bar = slide.shapes.add_shape(1, Inches(9.05), Inches(y + 0.03), Inches(max(bar_width, 0.05)), Inches(0.14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb("dark_blue")
            bar.line.fill.background()
            value_box = slide.shapes.add_textbox(Inches(12.02), Inches(y - 0.01), Inches(0.45), Inches(0.18))
            value_box.text_frame.text = str(int(value))
            value_box.text_frame.paragraphs[0].font.size = Pt(7)
            value_box.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_purple")
            y += 0.42

    def _add_pptx_leadership_brief(self, prs: Presentation, report: AccountReport, source_numbers: dict[str, int]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb("dark_blue")
        accent = slide.shapes.add_shape(1, Inches(0), Inches(6.72), Inches(13.333), Inches(0.22))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb("tech_purple")
        accent.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.58), Inches(0.42), Inches(9.6), Inches(0.55))
        title.text_frame.text = "Leadership Brief: So What"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.color.rgb = _rgb("light_blue")
        subtitle = slide.shapes.add_textbox(Inches(0.62), Inches(1.0), Inches(11.8), Inches(0.35))
        subtitle.text_frame.text = "The export now turns the evidence graph into account moves, not just section summaries."
        subtitle.text_frame.paragraphs[0].font.size = Pt(11)
        subtitle.text_frame.paragraphs[0].font.color.rgb = _rgb("tech_gray")

        x_positions = [0.6, 4.55, 8.5]
        for idx, column in enumerate(_leadership_brief(report)):
            rows: list[EvidenceTableRow] = column["rows"]  # type: ignore[assignment]
            panel = slide.shapes.add_shape(1, Inches(x_positions[idx]), Inches(1.55), Inches(3.55), Inches(4.85))
            panel.fill.solid()
            panel.fill.fore_color.rgb = _rgb("light_blue")
            panel.line.color.rgb = _rgb("tech_blue")
            heading = slide.shapes.add_textbox(Inches(x_positions[idx] + 0.18), Inches(1.78), Inches(3.15), Inches(0.35))
            heading.text_frame.text = str(column["title"])
            heading.text_frame.paragraphs[0].font.bold = True
            heading.text_frame.paragraphs[0].font.size = Pt(14)
            heading.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_blue")
            note = slide.shapes.add_textbox(Inches(x_positions[idx] + 0.18), Inches(2.17), Inches(3.1), Inches(0.45))
            note_frame = note.text_frame
            note_frame.word_wrap = True
            note_frame.text = str(column["subtitle"])
            note_frame.paragraphs[0].font.size = Pt(7.2)
            note_frame.paragraphs[0].font.color.rgb = _rgb("dark_purple")
            y = 2.68
            for row in rows:
                row_box = slide.shapes.add_textbox(Inches(x_positions[idx] + 0.18), Inches(y), Inches(3.12), Inches(1.06))
                row_frame = row_box.text_frame
                row_frame.word_wrap = True
                row_frame.margin_left = Inches(0.02)
                row_frame.margin_right = Inches(0.02)
                row_frame.margin_top = Inches(0.01)
                row_frame.margin_bottom = Inches(0.01)
                row_frame.clear()
                p = row_frame.paragraphs[0]
                p.text = _row_display_text(row, source_numbers)
                p.font.bold = True
                p.font.size = Pt(6.4)
                p.font.color.rgb = _rgb("dark_blue")
                detail = row_frame.add_paragraph()
                detail.text = f"{_row_strength(row)} | {_source_refs(row.source_ids, source_numbers, limit=3)}"
                detail.font.size = Pt(6.0)
                detail.font.color.rgb = _rgb("dark_purple")
                y += 1.18

        footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.25))
        footer.text_frame.text = "HCLTech Market Research Portal | Executive synthesis"
        footer.text_frame.paragraphs[0].font.size = Pt(8)
        footer.text_frame.paragraphs[0].font.color.rgb = _rgb("light_blue")
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _add_pptx_leadership_details(self, prs: Presentation, report: AccountReport, source_numbers: dict[str, int]) -> None:
        for column in _leadership_brief(report):
            rows: list[EvidenceTableRow] = column["rows"]  # type: ignore[assignment]
            if not rows:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = _rgb("light_blue")
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.82))
            header.fill.solid()
            header.fill.fore_color.rgb = _rgb("dark_blue")
            header.line.fill.background()
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(10.8), Inches(0.45))
            title_box.text_frame.text = f"Leadership Brief Details - {column['title']}"
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.size = Pt(21)
            title_box.text_frame.paragraphs[0].font.color.rgb = _rgb("tech_gray")
            subtitle = slide.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(12.0), Inches(0.32))
            subtitle.text_frame.text = str(column["subtitle"])
            subtitle.text_frame.paragraphs[0].font.size = Pt(9)
            subtitle.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_purple")

            y = 1.5
            for idx, row in enumerate(rows, start=1):
                box = slide.shapes.add_shape(1, Inches(0.62), Inches(y), Inches(12.05), Inches(1.48))
                box.fill.solid()
                box.fill.fore_color.rgb = _rgb("tech_gray" if idx % 2 else "light_blue")
                box.line.color.rgb = _rgb("mid_blue")
                text = slide.shapes.add_textbox(Inches(0.84), Inches(y + 0.12), Inches(11.55), Inches(1.22))
                frame = text.text_frame
                frame.word_wrap = True
                frame.margin_left = Inches(0.02)
                frame.margin_right = Inches(0.02)
                frame.margin_top = Inches(0.01)
                frame.margin_bottom = Inches(0.01)
                frame.clear()
                p = frame.paragraphs[0]
                p.text = f"{idx}. {_row_display_text(row, source_numbers)}"
                p.font.bold = True
                p.font.size = Pt(10)
                p.font.color.rgb = _rgb("dark_blue")
                detail = frame.add_paragraph()
                detail.text = f"{_row_strength(row)} | {_source_refs(row.source_ids, source_numbers, limit=8)}"
                detail.font.size = Pt(8)
                detail.font.color.rgb = _rgb("dark_purple")
                y += 1.66

            footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.25))
            footer.text_frame.text = "HCLTech Market Research Portal | Leadership detail"
            footer.text_frame.paragraphs[0].font.size = Pt(8)
            footer.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_blue")
            footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _add_pptx_insight_inventory(self, prs: Presentation, report: AccountReport, source_numbers: dict[str, int]) -> None:
        for section in _insight_inventory(report, limit_per_section=7):
            rows: list[EvidenceTableRow] = section["rows"]  # type: ignore[assignment]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = _rgb("light_blue")
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.82))
            header.fill.solid()
            header.fill.fore_color.rgb = _rgb("dark_blue")
            header.line.fill.background()
            accent = slide.shapes.add_shape(1, Inches(0), Inches(0.82), Inches(13.333), Inches(0.06))
            accent.fill.solid()
            accent.fill.fore_color.rgb = _rgb("tech_purple")
            accent.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.19), Inches(9.6), Inches(0.45))
            title_box.text_frame.text = str(section["title"])
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.size = Pt(22)
            title_box.text_frame.paragraphs[0].font.color.rgb = _rgb("tech_gray")

            count_box = slide.shapes.add_textbox(Inches(10.45), Inches(0.22), Inches(2.35), Inches(0.28))
            count_box.text_frame.text = f"{len(rows)} surfaced rows"
            count_box.text_frame.paragraphs[0].font.size = Pt(8)
            count_box.text_frame.paragraphs[0].font.color.rgb = _rgb("light_blue")
            count_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            subtitle = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(11.8), Inches(0.35))
            subtitle.text_frame.text = _truncate(str(section["subtitle"]), 165)
            subtitle.text_frame.paragraphs[0].font.size = Pt(10)
            subtitle.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_purple")

            table_x = 0.58
            table_y = 1.55
            widths = [2.7, 6.2, 1.35, 1.65]
            headers = ["Insight", "Evidence basis", "Strength", "Sources"]
            x = table_x
            for idx, header_text in enumerate(headers):
                box = slide.shapes.add_shape(1, Inches(x), Inches(table_y), Inches(widths[idx]), Inches(0.32))
                box.fill.solid()
                box.fill.fore_color.rgb = _rgb("dark_blue")
                box.line.fill.background()
                frame = box.text_frame
                frame.clear()
                p = frame.paragraphs[0]
                p.text = header_text
                p.font.bold = True
                p.font.size = Pt(8)
                p.font.color.rgb = _rgb("light_blue")
                x += widths[idx]

            row_y = table_y + 0.38
            row_h = 0.66
            for idx, row in enumerate(rows):
                fill = "tech_gray" if idx % 2 == 0 else "light_blue"
                x = table_x
                values = [
                    _truncate(_clean_text(row.title, source_numbers), 78),
                    _truncate(_clean_text(row.detail, source_numbers), 210),
                    _truncate(_row_strength(row), 34),
                    _source_refs(row.source_ids, source_numbers),
                ]
                for col, value in enumerate(values):
                    cell = slide.shapes.add_shape(1, Inches(x), Inches(row_y), Inches(widths[col]), Inches(row_h))
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(fill)
                    cell.line.color.rgb = _rgb("mid_blue")
                    frame = cell.text_frame
                    frame.word_wrap = True
                    frame.margin_left = Inches(0.05)
                    frame.margin_right = Inches(0.05)
                    frame.margin_top = Inches(0.03)
                    frame.margin_bottom = Inches(0.03)
                    frame.clear()
                    p = frame.paragraphs[0]
                    p.text = value
                    p.font.size = Pt(6.5 if col == 1 else 7)
                    p.font.bold = col == 0
                    p.font.color.rgb = _rgb("dark_purple")
                    x += widths[col]
                row_y += row_h + 0.04

            footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.25))
            footer.text_frame.text = "HCLTech Market Research Portal | Insight inventory from evidence tables"
            footer.text_frame.paragraphs[0].font.size = Pt(8)
            footer.text_frame.paragraphs[0].font.color.rgb = _rgb("dark_blue")
            footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def export_pdf(self, report: AccountReport, path: Path) -> Artifact:
        source_numbers = _source_numbers(report)
        page = canvas.Canvas(str(path), pagesize=landscape(letter))
        width, height = landscape(letter)
        page_number = 0

        def footer() -> None:
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica", 7)
            page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

        def new_content_page(title: str, subtitle: str | None = None) -> None:
            nonlocal page_number
            if page_number:
                page.showPage()
            page_number += 1
            page.setFillColor(colors.white)
            page.rect(0, 0, width, height, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.rect(0, height - 58, width, 58, fill=1, stroke=0)
            page.setFillColor(_rl_color("tech_purple"))
            page.rect(0, height - 64, width, 6, fill=1, stroke=0)
            page.setFillColor(colors.white)
            page.setFont("Helvetica-Bold", 19)
            page.drawString(34, height - 36, title[:95])
            if subtitle:
                page.setFont("Helvetica", 8)
                page.drawRightString(width - 34, height - 35, subtitle[:95])
            footer()

        page_number += 1
        page.setFillColor(_rl_color("dark_blue"))
        page.rect(0, 0, width, height, fill=1, stroke=0)
        page.setFillColor(_rl_color("tech_purple"))
        page.rect(0, 72, width, 18, fill=1, stroke=0)
        page.setFillColor(_rl_color("light_blue"))
        page.setFont("Helvetica-Bold", 30)
        page.drawString(52, height - 150, f"{report.company_name} Account Intelligence")
        page.setFont("Helvetica", 13)
        page.drawString(54, height - 181, "Source-grounded account intelligence, AI strategy, and HCLTech pursuit guidance")
        page.setFont("Helvetica", 10)
        metadata = f"Mode: {report.mode.value.title()} | Freshness: {report.freshness_window.value} | Claims: {len(report.claims)} | Public sources: {len(_public_sources(report))}"
        page.drawString(54, height - 224, metadata)
        card_x = 54
        for label, value in _metric_counts(report).items():
            self._draw_pdf_metric_card(page, card_x, height - 306, label, value)
            card_x += 132
        page.setFillColor(colors.white)
        page.setFont("Helvetica", 8)
        page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

        page_number = self._draw_pdf_summary_pages(page, report, source_numbers, width, height, page_number)
        page_number = self._draw_pdf_leadership_brief(page, report, source_numbers, width, height, page_number)
        page_number = self._draw_pdf_leadership_details(page, report, source_numbers, width, height, page_number)
        page_number = self._draw_pdf_insight_inventory(page, report, source_numbers, width, height, page_number)

        for section in report.sections:
            new_content_page(section.title, f"Status: {section.status} | Confidence: {section.confidence_score:.2f}")
            claims = _claims_for_section(report, section)
            citation_labels = _citation_labels(report, _citation_ids_for_section(report, section), source_numbers, limit=7)
            page.setFillColor(_rl_color("tech_gray"))
            page.roundRect(34, height - 92, 494, 24, 5, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 8)
            page.drawString(48, height - 78, f"{section.status.upper()} | Confidence {section.confidence_score:.2f}")
            page.setFillColor(_rl_color("dark_purple"))
            page.setFont("Helvetica", 8)
            page.drawRightString(512, height - 78, f"{len(claims)} mapped claims | {len(citation_labels)} displayed sources")

            y = height - 112
            page.setFillColor(_rl_color("light_blue"))
            page.roundRect(34, y - 114, width - 298, 112, 8, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_purple"))
            page.setFont("Helvetica-Bold", 11)
            page.drawString(52, y - 24, "Synthesis")
            y = _draw_wrapped(
                page,
                _clean_text(section.summary, source_numbers),
                52,
                y - 44,
                width - 334,
                "Helvetica",
                9,
                12,
                _rl_color("dark_purple"),
                max_lines=6,
            )

            y -= 14
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 10)
            page.drawString(52, y, "Claim-level support")
            y -= 17
            for claim in claims[:4]:
                claim_text = f"- {_clean_text(claim.text, source_numbers)}"
                y = _draw_wrapped(page, claim_text, 58, y, width - 342, "Helvetica", 8, 10, colors.black, max_lines=3)
                y -= 5
                if y < 76:
                    page.setFillColor(_rl_color("dark_blue"))
                    page.setFont("Helvetica-Oblique", 8)
                    page.drawString(58, y, "Additional claim detail is available in the evidence pack.")
                    break

            panel_x = width - 238
            panel_y = height - 96
            page.setFillColor(_rl_color("tech_gray"))
            page.roundRect(panel_x, 74, 204, panel_y - 76, 8, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 10)
            page.drawString(panel_x + 14, panel_y - 24, "Evidence")
            cite_y = panel_y - 44
            for label in citation_labels or ["No public source attached."]:
                cite_y = _draw_wrapped(page, label, panel_x + 14, cite_y, 176, "Helvetica", 7, 9, _rl_color("dark_purple"), max_lines=3)
                cite_y -= 7

        page_number = self._draw_sources_appendix(page, report, source_numbers, width, height, page_number)
        page_number = self._draw_claim_appendix(page, report, source_numbers, width, height, page_number)
        page.save()

        raw_ids_in_reader_text = any(SOURCE_ID_RE.search(_clean_text(section.summary, source_numbers)) for section in report.sections)
        scaffold_terms = [
            section.title
            for section in report.sections
            if any(term in _clean_text(section.summary, source_numbers).lower() for term in SCAFFOLD_TERMS)
        ]
        return Artifact(
            kind="pdf",
            path=str(path),
            quality_checks=[
                QualityCheck(name="pdf_created", passed=path.exists(), message=f"PDF written to {path}"),
                QualityCheck(name="pdf_uses_brand_palette", passed=True, message="PDF pages use mandatory palette tokens."),
                QualityCheck(
                    name="pdf_human_readable_citations",
                    passed=not raw_ids_in_reader_text,
                    severity="blocker" if raw_ids_in_reader_text else "info",
                    message="PDF replaces internal source IDs with S-number citation labels.",
                ),
                QualityCheck(
                    name="pdf_reader_ready_content",
                    passed=not scaffold_terms,
                    severity="warning" if scaffold_terms else "info",
                    message=(
                        f"Scaffold language remains in sections: {', '.join(scaffold_terms[:5])}."
                        if scaffold_terms
                        else "No scaffold language detected in reader-facing summaries."
                    ),
                ),
            ],
        )

    def _draw_pdf_header(self, page: canvas.Canvas, title: str, width: float, height: float, page_number: int, subtitle: str | None = None) -> None:
        page.setFillColor(colors.white)
        page.rect(0, 0, width, height, fill=1, stroke=0)
        page.setFillColor(_rl_color("dark_blue"))
        page.rect(0, height - 58, width, 58, fill=1, stroke=0)
        page.setFillColor(_rl_color("tech_purple"))
        page.rect(0, height - 64, width, 6, fill=1, stroke=0)
        page.setFillColor(colors.white)
        page.setFont("Helvetica-Bold", 19)
        page.drawString(34, height - 36, title[:95])
        if subtitle:
            page.setFont("Helvetica", 8)
            page.drawRightString(width - 34, height - 35, subtitle[:95])
        page.setFillColor(_rl_color("dark_blue"))
        page.setFont("Helvetica", 7)
        page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

    def _draw_pdf_metric_card(self, page: canvas.Canvas, x: float, y: float, label: str, value: int) -> None:
        page.setFillColor(_rl_color("light_blue"))
        page.roundRect(x, y - 58, 118, 58, 7, fill=1, stroke=0)
        page.setFillColor(_rl_color("dark_blue"))
        page.setFont("Helvetica-Bold", 20)
        page.drawString(x + 14, y - 26, str(value))
        page.setFillColor(_rl_color("dark_purple"))
        page.setFont("Helvetica", 8)
        page.drawString(x + 14, y - 44, label)

    def _draw_pdf_bar_chart(self, page: canvas.Canvas, title: str, data: list[dict], x: float, y: float, width: float, max_items: int = 8) -> float:
        page.setFillColor(_rl_color("dark_blue"))
        page.setFont("Helvetica-Bold", 12)
        page.drawString(x, y, title)
        y -= 22
        if not data:
            return y
        max_value = max(float(item.get("value", 0) or 0) for item in data[:max_items]) or 1
        for item in data[:max_items]:
            label = str(item.get("label", ""))[:42]
            value = float(item.get("value", 0) or 0)
            page.setFillColor(_rl_color("dark_purple"))
            page.setFont("Helvetica", 7)
            page.drawString(x, y, label)
            bar_x = x + 190
            bar_w = max(2, (width - 235) * value / max_value)
            page.setFillColor(_rl_color("light_blue"))
            page.roundRect(bar_x, y - 2, width - 235, 8, 3, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.roundRect(bar_x, y - 2, bar_w, 8, 3, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_purple"))
            page.setFont("Helvetica-Bold", 7)
            page.drawRightString(x + width, y, str(int(value)))
            y -= 20
        return y

    def _draw_pdf_summary_pages(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        page.showPage()
        page_number += 1
        self._draw_pdf_header(page, "Executive Readout", width, height, page_number, "Evidence quality before section detail")
        metrics = _metric_counts(report)
        x = 46
        for label, value in metrics.items():
            self._draw_pdf_metric_card(page, x, height - 110, label, value)
            x += 134

        verified = len([claim for claim in report.claims if claim.verification_status == "verified"])
        unavailable = len([claim for claim in report.claims if claim.verification_status == "unavailable"])
        weak_sections = [section.title for section in report.sections if section.status == "unavailable" or section.confidence_score < 0.6]
        strong_sections = [section.title for section in report.sections if section.status != "unavailable" and section.confidence_score >= 0.7]
        y = height - 195
        y = _draw_wrapped(
            page,
            f"This report currently has {verified} verified claims and {unavailable} unavailable claims. Unavailable claims are retained deliberately so unsupported market intelligence is not invented.",
            46,
            y,
            width - 92,
            "Helvetica",
            10,
            13,
            _rl_color("dark_purple"),
            max_lines=3,
            add_ellipsis=False,
        )
        y -= 16
        y = _draw_wrapped(
            page,
            f"Strongest evidence coverage: {', '.join(strong_sections[:5]) or 'still developing'}.",
            46,
            y,
            width - 92,
            "Helvetica-Bold",
            10,
            13,
            _rl_color("dark_blue"),
            max_lines=3,
        )
        y -= 10
        _draw_wrapped(
            page,
            f"Open gaps before client-ready use: {', '.join(weak_sections[:6]) or 'no major gaps flagged'}.",
            46,
            y,
            width - 92,
            "Helvetica",
            9,
            12,
            _rl_color("dark_purple"),
            max_lines=4,
        )

        page.showPage()
        page_number += 1
        self._draw_pdf_header(page, "Evidence Coverage", width, height, page_number, "Charts generated from traceability metadata")
        self._draw_pdf_bar_chart(page, "Claim Status", _claim_status_data(report), 46, height - 100, 310, max_items=4)
        self._draw_pdf_bar_chart(page, "Public Sources by Tier", _source_mix_data(report), 420, height - 100, 320, max_items=8)
        self._draw_pdf_bar_chart(page, "Evidence Signals by Type", _signal_mix_data(report), 46, height - 300, 310, max_items=8)
        self._draw_pdf_bar_chart(page, "Evidence Rows by Table", _table_mix_data(report), 420, height - 300, 320, max_items=8)
        return page_number

    def _draw_pdf_leadership_brief(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        page.showPage()
        page_number += 1
        page.setFillColor(_rl_color("dark_blue"))
        page.rect(0, 0, width, height, fill=1, stroke=0)
        page.setFillColor(_rl_color("tech_purple"))
        page.rect(0, height - 78, width, 8, fill=1, stroke=0)
        page.setFillColor(_rl_color("light_blue"))
        page.setFont("Helvetica-Bold", 24)
        page.drawString(42, height - 48, "Leadership Brief: So What")
        page.setFont("Helvetica", 9)
        page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

        x_positions = [42, 292, 542]
        panel_w = 216
        for idx, column in enumerate(_leadership_brief(report)):
            rows: list[EvidenceTableRow] = column["rows"]  # type: ignore[assignment]
            x = x_positions[idx]
            page.setFillColor(_rl_color("light_blue"))
            page.roundRect(x, height - 430, panel_w, 315, 8, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 13)
            page.drawString(x + 14, height - 140, str(column["title"]))
            _draw_wrapped(
                page,
                str(column["subtitle"]),
                x + 14,
                height - 160,
                panel_w - 28,
                "Helvetica",
                7,
                9,
                _rl_color("dark_purple"),
                max_lines=3,
            )
            y = height - 205
            for row in rows:
                _draw_wrapped(
                    page,
                    _row_display_text(row, source_numbers),
                    x + 14,
                    y,
                    panel_w - 28,
                    "Helvetica-Bold",
                    7,
                    8,
                    _rl_color("dark_blue"),
                    max_lines=4,
                    add_ellipsis=False,
                )
                _draw_wrapped(
                    page,
                    f"{_row_strength(row)} | {_source_refs(row.source_ids, source_numbers, limit=3)}",
                    x + 14,
                    y - 32,
                    panel_w - 28,
                    "Helvetica",
                    6,
                    7,
                    _rl_color("dark_purple"),
                    max_lines=1,
                    add_ellipsis=False,
                )
                y -= 74
        return page_number

    def _draw_pdf_leadership_details(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        for column in _leadership_brief(report):
            rows: list[EvidenceTableRow] = column["rows"]  # type: ignore[assignment]
            if not rows:
                continue
            page.showPage()
            page_number += 1
            self._draw_pdf_header(page, f"Leadership Brief Details - {column['title']}", width, height, page_number, str(column["subtitle"]))
            y = height - 100
            for idx, row in enumerate(rows, start=1):
                box_top = y + 10
                box_h = 102
                page.setFillColor(_rl_color("tech_gray" if idx % 2 else "light_blue"))
                page.roundRect(40, box_top - box_h, width - 80, box_h, 6, fill=1, stroke=0)
                y = _draw_wrapped(
                    page,
                    f"{idx}. {_row_display_text(row, source_numbers)}",
                    54,
                    box_top - 22,
                    width - 108,
                    "Helvetica-Bold",
                    9,
                    11,
                    _rl_color("dark_blue"),
                    max_lines=5,
                    add_ellipsis=False,
                )
                y = _draw_wrapped(
                    page,
                    f"{_row_strength(row)} | {_source_refs(row.source_ids, source_numbers, limit=8)}",
                    54,
                    y - 6,
                    width - 108,
                    "Helvetica",
                    7,
                    9,
                    _rl_color("dark_purple"),
                    max_lines=2,
                    add_ellipsis=False,
                )
                y = box_top - box_h - 16
        return page_number

    def _draw_pdf_insight_inventory(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        inventory = _insight_inventory(report, limit_per_section=7)
        for section in inventory:
            rows: list[EvidenceTableRow] = section["rows"]  # type: ignore[assignment]
            page.showPage()
            page_number += 1
            self._draw_pdf_header(page, str(section["title"]), width, height, page_number, f"{len(rows)} surfaced insight rows")
            y = height - 88
            y = _draw_wrapped(
                page,
                str(section["subtitle"]),
                46,
                y,
                width - 92,
                "Helvetica",
                8,
                10,
                _rl_color("dark_purple"),
                max_lines=2,
            )
            y -= 12
            columns = [
                ("Insight", 46, 160),
                ("Evidence basis", 210, 344),
                ("Strength", 560, 76),
                ("Sources", 642, 100),
            ]
            page.setFillColor(_rl_color("dark_blue"))
            page.roundRect(34, y - 20, width - 68, 20, 4, fill=1, stroke=0)
            page.setFillColor(_rl_color("light_blue"))
            page.setFont("Helvetica-Bold", 7)
            for label, x, _ in columns:
                page.drawString(x, y - 13, label)
            y -= 30
            for idx, row in enumerate(rows):
                if y < 86:
                    page.showPage()
                    page_number += 1
                    self._draw_pdf_header(page, str(section["title"]), width, height, page_number, "continued")
                    y = height - 88
                row_height = 50
                page.setFillColor(_rl_color("tech_gray" if idx % 2 == 0 else "light_blue"))
                page.roundRect(34, y - row_height + 4, width - 68, row_height, 4, fill=1, stroke=0)
                page.setFillColor(_rl_color("dark_blue"))
                page.setFont("Helvetica-Bold", 7)
                _draw_wrapped(
                    page,
                    _clean_text(row.title, source_numbers),
                    46,
                    y - 10,
                    150,
                    "Helvetica-Bold",
                    7,
                    8,
                    _rl_color("dark_blue"),
                    max_lines=3,
                )
                _draw_wrapped(
                    page,
                    _clean_text(row.detail, source_numbers),
                    210,
                    y - 10,
                    336,
                    "Helvetica",
                    7,
                    8,
                    colors.black,
                    max_lines=4,
                )
                _draw_wrapped(
                    page,
                    _row_strength(row),
                    560,
                    y - 10,
                    70,
                    "Helvetica-Bold",
                    7,
                    8,
                    _rl_color("dark_purple"),
                    max_lines=2,
                )
                _draw_wrapped(
                    page,
                    _source_refs(row.source_ids, source_numbers),
                    642,
                    y - 10,
                    92,
                    "Helvetica",
                    7,
                    8,
                    _rl_color("dark_purple"),
                    max_lines=2,
                )
                y -= row_height + 6
        return page_number

    def _draw_sources_appendix(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        sources = _public_sources(report)
        if not sources:
            return page_number
        page.showPage()
        page_number += 1
        page.setFillColor(colors.white)
        page.rect(0, 0, width, height, fill=1, stroke=0)
        page.setFillColor(_rl_color("dark_blue"))
        page.rect(0, height - 58, width, 58, fill=1, stroke=0)
        page.setFillColor(colors.white)
        page.setFont("Helvetica-Bold", 19)
        page.drawString(34, height - 36, "Sources Appendix")
        page.setFillColor(_rl_color("dark_blue"))
        page.setFont("Helvetica", 7)
        page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

        y = height - 88
        for source in sources[:32]:
            if y < 74:
                page.showPage()
                page_number += 1
                page.setFillColor(colors.white)
                page.rect(0, 0, width, height, fill=1, stroke=0)
                page.setFillColor(_rl_color("dark_blue"))
                page.rect(0, height - 58, width, 58, fill=1, stroke=0)
                page.setFillColor(colors.white)
                page.setFont("Helvetica-Bold", 19)
                page.drawString(34, height - 36, "Sources Appendix")
                page.setFillColor(_rl_color("dark_blue"))
                page.setFont("Helvetica", 7)
                page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")
                y = height - 88
            number = source_numbers.get(source.id, 0)
            title = f"S{number}: {_readable_source_title(source)}"
            publisher = f"{_safe_text(source.publisher or 'Unknown publisher')} | {source.credibility.value} | score {source.credibility_score:.2f}"
            page.setFillColor(_rl_color("light_blue"))
            page.roundRect(34, y - 46, width - 68, 42, 5, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 8)
            page.drawString(46, y - 17, title[:132])
            page.setFillColor(_rl_color("dark_purple"))
            page.setFont("Helvetica", 7)
            page.drawString(46, y - 29, publisher[:150])
            page.setFillColor(colors.black)
            page.setFont("Helvetica", 6)
            page.drawString(46, y - 39, _safe_text(source.url[:170]))
            y -= 52
        return page_number

    def _draw_claim_appendix(
        self,
        page: canvas.Canvas,
        report: AccountReport,
        source_numbers: dict[str, int],
        width: float,
        height: float,
        page_number: int,
    ) -> int:
        page.showPage()
        page_number += 1
        page.setFillColor(colors.white)
        page.rect(0, 0, width, height, fill=1, stroke=0)
        page.setFillColor(_rl_color("dark_blue"))
        page.rect(0, height - 58, width, 58, fill=1, stroke=0)
        page.setFillColor(colors.white)
        page.setFont("Helvetica-Bold", 19)
        page.drawString(34, height - 36, "Claim Evidence Matrix")
        page.setFillColor(_rl_color("dark_blue"))
        page.setFont("Helvetica", 7)
        page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")

        y = height - 86
        for claim in report.claims[:28]:
            if y < 86:
                page.showPage()
                page_number += 1
                page.setFillColor(colors.white)
                page.rect(0, 0, width, height, fill=1, stroke=0)
                page.setFillColor(_rl_color("dark_blue"))
                page.rect(0, height - 58, width, 58, fill=1, stroke=0)
                page.setFillColor(colors.white)
                page.setFont("Helvetica-Bold", 19)
                page.drawString(34, height - 36, "Claim Evidence Matrix")
                page.setFillColor(_rl_color("dark_blue"))
                page.setFont("Helvetica", 7)
                page.drawRightString(width - 34, 22, f"HCLTech Market Research Portal | {page_number}")
                y = height - 86
            evidence_labels = [
                f"S{source_numbers[source_id]}"
                for source_id in claim.evidence_source_ids
                if source_id in source_numbers
            ]
            status = f"{claim.verification_status} | {claim.claim_type} | {', '.join(evidence_labels) or 'internal/system evidence'}"
            page.setFillColor(_rl_color("tech_gray"))
            page.roundRect(34, y - 45, width - 68, 40, 5, fill=1, stroke=0)
            page.setFillColor(_rl_color("dark_blue"))
            page.setFont("Helvetica-Bold", 7)
            page.drawString(46, y - 17, status[:150])
            _draw_wrapped(
                page,
                _clean_text(claim.text, source_numbers),
                46,
                y - 30,
                width - 92,
                "Helvetica",
                7,
                8,
                colors.black,
                max_lines=2,
            )
            y -= 48
        return page_number

    def export_evidence(self, report: AccountReport, path: Path) -> Artifact:
        payload = {
            "run_id": report.run_id,
            "company_name": report.company_name,
            "claims": [claim.model_dump(mode="json") for claim in report.claims],
            "sources": [source.model_dump(mode="json") for source in report.sources],
            "snapshots": [snapshot.model_dump(mode="json") for snapshot in report.snapshots],
            "extracted_values": [value.model_dump(mode="json") for value in report.extracted_values],
            "evidence_signals": [signal.model_dump(mode="json") for signal in report.evidence_signals],
            "evidence_table_rows": [row.model_dump(mode="json") for row in report.evidence_table_rows],
            "sections": [section.model_dump(mode="json") for section in report.sections],
            "quality_checks": [check.model_dump(mode="json") for check in report.quality_checks],
        }
        path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        return Artifact(
            kind="evidence_json",
            path=str(path),
            quality_checks=[QualityCheck(name="evidence_pack_created", passed=path.exists(), message=f"Evidence pack written to {path}")],
        )
